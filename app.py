# 1. Importación de módulos de Flask para funcionalidad web básica
from flask import Flask, render_template, request, send_file, url_for

# 2. Importación de Google Generative AI para interactuar con Gemini
import google.generativeai as genai

# 3. Importaciones para manejo de documentos Word
from docx import Document
from docx.shared import Pt, RGBColor

# 4. Importaciones para generación de PDFs
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# 5. Importaciones de bibliotecas estándar de Python
import os
import io

# 6. Importación de Pandas para manejo de datos tabulares
import pandas as pd

# 7. Importaciones para manejo de Excel
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Side, Border

# 8. Importaciones adicionales para Word
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches
from docx.oxml import parse_xml

# 9. Importaciones para PowerPoint
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor as PPTXRGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# 10. Inicialización de la aplicación Flask
app = Flask(__name__, static_folder='static')

# 11. Configuración de credenciales para Gemini
genai.configure(api_key='AQUI TU API KEY')
model = genai.GenerativeModel('gemini-1.5-flash')

# 12. Función auxiliar para procesar contenido Excel/CSV
def parse_excel_content(content):
    """
    Función para parsear contenido en formato Excel/CSV
    Entrada: String con contenido separado por comas y saltos de línea
    Salida: Lista de listas donde cada sublista es una fila de datos
    """
    rows = []
    for line in content.strip().split('\n'):
        # Divide cada línea por comas y limpia espacios
        row = [cell.strip() for cell in line.split(',')]
        rows.append(row)
    return rows

# 13. Definición de la ruta principal
@app.route('/', methods=['GET', 'POST'])
def index():
    """
    Ruta principal que maneja:
    GET: Muestra el formulario inicial
    POST: Procesa la solicitud y genera el documento
    """
    if request.method == 'POST':
        try:
            # 14. Obtención de datos del formulario
            prompt = request.form['prompt']  # Texto ingresado por usuario
            formato = request.form['formato']  # Formato de salida deseado
            
            # 15. Lógica para generación de presentación PowerPoint
            if formato == 'pptx':
                # Template para estructurar la presentación
                prompt_template = f"""Crea una presentación de PowerPoint sobre: {prompt}

                Estructura exacta requerida (mantén este formato):
                [Diapositiva 1]
                Título: [título corto y claro]
                Subtítulo: [subtítulo descriptivo]

                [Diapositiva 2]
                Título: Introducción
                • [punto clave 1]
                • [punto clave 2]
                • [punto clave 3]

                [Diapositiva 3]
                Título: Desarrollo
                • [punto principal 1]
                • [detalle importante]
                • [ejemplo concreto]

                [Diapositiva 4]
                Título: Conclusiones
                • [conclusión principal]
                • [recomendación]
                • [punto final]

                Usa lenguaje claro y directo. Cada punto debe ser una frase completa."""

                # 16. Generación de contenido con Gemini
                response = model.generate_content(prompt_template)
                
                # 17. Validación de respuesta
                if not response or not response.text:
                    return "Error: No se pudo generar el contenido. Por favor, intente de nuevo."

                print("Contenido generado:", response.text)  # Debug

            # 18. Lógica para generación de Excel/CSV    
            elif formato in ['xlsx', 'csv']:
                # Genera contenido tabular
                response = model.generate_content([
                    """Genera una estructura de datos tabulada para Excel sobre el tema.
                    Debes generar:
                    1. Una primera línea con los nombres de las columnas
                    2. Al menos 10 filas de datos de ejemplo relevantes al tema
                    3. Los datos deben estar separados por comas
                    
                    Formato exacto:
                    columna1,columna2,columna3,columna4
                    dato1,dato2,dato3,dato4
                    dato1,dato2,dato3,dato4
                    
                    No incluyas ningún mensaje adicional.
                    
                    Tema:""",
                    prompt
                ])

            # 19. Lógica para generación de otros formatos
            else:
                # Template para documentos de texto
                response = model.generate_content([
                    """Basado en el siguiente prompt, genera:
                    1. Un título apropiado y profesional (en una sola línea)
                    2. Desarrolla el contenido en formato de texto estructurado con:
                       - Introducción
                       - Desarrollo de puntos principales
                       - Conclusión
                    
                    Reglas:
                    - El título debe estar en la primera línea
                    - Dejar una línea en blanco después del título
                    - No usar formato de tabla
                    - Usar párrafos bien estructurados
                    - Incluir subtítulos relevantes
                    - Mantener un tono profesional
                    
                    Formato esperado:
                    [Título del documento]

                    Introducción:
                    [Contenido]
                    
                    [Subtítulo 1]:
                    [Contenido]
                    
                    [Subtítulo 2]:
                    [Contenido]
                    
                    Conclusión:
                    [Contenido]
                    
                    Prompt:""",
                    prompt
                ])
            
            # 20. Procesamiento específico para PowerPoint
            if formato == 'pptx':
                try:
                    # 21. Crear nueva presentación
                    prs = Presentation()
                    
                    # 22. Definición de estilos consistentes
                    title_font_size = Pt(44)
                    subtitle_font_size = Pt(32)
                    body_font_size = Pt(24)
                    
                    # 23. Colores para el texto
                    title_color = PPTXRGBColor(255, 255, 255)  # Blanco
                    body_color = PPTXRGBColor(255, 255, 255)  # Blanco
                    
                    # 24. Función para aplicar fondo a diapositivas
                    def aplicar_fondo(slide):
                        """
                        Aplica un fondo degradado oscuro a la diapositiva
                        """
                        left = top = 0
                        width = prs.slide_width
                        height = prs.slide_height
                        
                        # Crear rectángulo de fondo
                        shape = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE, left, top, width, height
                        )
                        
                        # Configurar gradiente
                        fill = shape.fill
                        fill.gradient()
                        fill.gradient_angle = 45
                        
                        # Colores del gradiente
                        fill.gradient_stops[0].position = 0
                        fill.gradient_stops[0].color.rgb = PPTXRGBColor(40, 40, 40)
                        fill.gradient_stops[1].position = 1
                        fill.gradient_stops[1].color.rgb = PPTXRGBColor(0, 0, 0)
                        
                        shape.line.fill.background()
                        
                        # Mover al fondo
                        shape_element = shape._element
                        slide.shapes._spTree.remove(shape_element)
                        slide.shapes._spTree.insert(0, shape_element)
                    
                    # 25. Función para crear diapositiva de título
                    def crear_diapositiva_titulo(title_text, subtitle_text=""):
                        """
                        Crea una diapositiva de título con estilo personalizado
                        """
                        slide = prs.slides.add_slide(prs.slide_layouts[0])
                        aplicar_fondo(slide)
                        
                        # Configurar título
                        if slide.shapes.title:
                            title = slide.shapes.title
                            title.text = title_text
                            for paragraph in title.text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.CENTER
                                for run in paragraph.runs:
                                    run.font.size = title_font_size
                                    run.font.color.rgb = title_color
                                    run.font.bold = True
                        
                        # Configurar subtítulo
                        if len(slide.placeholders) > 1:
                            subtitle = slide.placeholders[1]
                            subtitle.text = subtitle_text
                            for paragraph in subtitle.text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.CENTER
                                for run in paragraph.runs:
                                    run.font.size = subtitle_font_size
                                    run.font.color.rgb = title_color
                                    run.font.bold = True
                        
                        return slide
                    
                    # 26. Función para crear diapositiva de contenido
                    def crear_diapositiva_contenido(title_text, bullet_points):
                        """
                        Crea una diapositiva de contenido con viñetas
                        """
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        aplicar_fondo(slide)
                        
                        # Configurar título
                        if slide.shapes.title:
                            title = slide.shapes.title
                            title.text = title_text
                            for paragraph in title.text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.CENTER
                                for run in paragraph.runs:
                                    run.font.size = Pt(36)
                                    run.font.color.rgb = title_color
                                    run.font.bold = True
                        
                        # Configurar contenido con viñetas
                        if len(slide.placeholders) > 1:
                            body = slide.placeholders[1]
                            tf = body.text_frame
                            tf.clear()
                            
                            for point in bullet_points:
                                p = tf.add_paragraph()
                                p.text = point
                                p.level = 0
                                p.alignment = PP_ALIGN.LEFT
                                for run in p.runs:
                                    run.font.size = body_font_size
                                    run.font.color.rgb = body_color
                                    run.font.bold = True
                        
                        return slide
                    
                    # 27. Procesar contenido y crear diapositivas
                    slides_content = response.text.split('[Diapositiva')
                    
                    for slide_idx, slide_content in enumerate(slides_content[1:], 1):
                        lines = [line.strip() for line in slide_content.split('\n') if line.strip()]
                        
                        title = ""
                        subtitle = ""
                        bullet_points = []
                        
                        # Extraer contenido de la diapositiva
                        for line in lines:
                            if 'Título:' in line:
                                title = line.replace('Título:', '').strip()
                            elif 'Subtítulo:' in line:
                                subtitle = line.replace('Subtítulo:', '').strip()
                            elif '•' in line:
                                bullet_points.append(line.replace('•', '').strip())
                        
                        # Crear diapositiva según tipo
                        if slide_idx == 1:
                            crear_diapositiva_titulo(title, subtitle)
                        else:
                            crear_diapositiva_contenido(title, bullet_points)
                    
                    # 28. Guardar presentación
                    pptx_buffer = io.BytesIO()
                    prs.save(pptx_buffer)
                    pptx_buffer.seek(0)
                    
                    filename = f"presentacion_{prompt[:30].lower().replace(' ', '_')}.pptx"
                    
                    return send_file(
                        pptx_buffer,
                        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                        as_attachment=True,
                        download_name=filename
                    )
                    
                except Exception as e:
                    return {
                        'message': f"Error al crear la presentación PowerPoint: {str(e)}",
                        'error': True
                    }

            # 29. Procesamiento específico para Excel
            if formato == 'xlsx':
                # Parsear contenido
                rows = parse_excel_content(response.text)
                
                # Generar nombre de archivo
                filename = f"{rows[0][0].lower().replace(' ', '_')[:30]}.xlsx"
                
                try:
                    # Crear nuevo libro Excel
                    wb = Workbook()
                    ws = wb.active
                    
                    # Definir estilos
                    header_font = Font(name='Calibri', bold=True, size=12, color='FFFFFF')
                    header_fill = PatternFill(start_color='1F497D', end_color='1F497D', fill_type='solid')
                    header_alignment = Alignment(horizontal='center', vertical='center')
                    
                    data_font = Font(name='Calibri', size=11)
                    data_alignment = Alignment(horizontal='left', vertical='center')
                    
                    # Escribir datos y aplicar estilos
                    for i, row in enumerate(rows, 1):
                        for j, value in enumerate(row, 1):
                            cell = ws.cell(row=i, column=j, value=value)
                            if i == 1:  # Encabezados
                                cell.font = header_font
                                cell.fill = header_fill
                                cell.alignment = header_alignment
                            else:  # Datos
                                cell.font = data_font
                                cell.alignment = data_alignment
                    
                    # Ajustar ancho de columnas
                    for column in ws.columns:
                        max_length = 0
                        column_letter = column[0].column_letter
                        for cell in column:
                            try:
                                if len(str(cell.value)) > max_length:
                                    max_length = len(str(cell.value))
                            except:
                                pass
                        adjusted_width = (max_length + 2)
                        ws.column_dimensions[column_letter].width = adjusted_width
                    
                    # Guardar archivo
                    excel_buffer = io.BytesIO()
                    wb.save(excel_buffer)
                    excel_buffer.seek(0)
                    
                    return send_file(
                        excel_buffer,
                        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                        as_attachment=True,
                        download_name=filename
                    )
                    
                except Exception as e:
                    return {
                        'message': f"Error al crear el documento de Excel: {str(e)}",
                        'error': True
                    }
            
            # 30. Procesamiento específico para CSV
            if formato == 'csv':
                # Parsear contenido
                rows = parse_excel_content(response.text)
                
                # Crear DataFrame
                df = pd.DataFrame(rows[1:], columns=rows[0])
                
                # Generar nombre de archivo
                filename = f"{rows[0][0].lower().replace(' ', '_')[:30]}.csv"
                
                # Crear buffer y guardar CSV
                csv_buffer = io.StringIO()
                df.to_csv(csv_buffer, index=False, encoding='utf-8')
                
                output = io.BytesIO()
                output.write(csv_buffer.getvalue().encode('utf-8'))
                output.seek(0)
                
                return send_file(
                    output,
                    mimetype='text/csv',
                    as_attachment=True,
                    download_name=filename
                )
            
            # 31. Procesamiento de documentos de texto
            # Limpiar y formatear contenido
            contenido = response.text.strip().replace('**', '').replace('#', '').replace('*', '')
            lineas = contenido.split('\n')
            
            # Acortar título si es muy largo
            if lineas and len(lineas[0]) > 50:
                palabras = lineas[0].split(':')
                if len(palabras) > 1:
                    lineas[0] = palabras[0] + ':'
                else:
                    lineas[0] = ' '.join(lineas[0].split()[:6]) + '...'
            
            # Acortar subtítulos largos
            for i in range(1, len(lineas)):
                if lineas[i].strip() and len(lineas[i]) > 40 and lineas[i].endswith(':'):
                    palabras = lineas[i].split(':')
                    if len(palabras) > 1:
                        lineas[i] = palabras[0] + ':'
                    else:
                        lineas[i] = ' '.join(lineas[i].split()[:4]) + '...'
            
            contenido = '\n'.join(lineas)
            
            # 32. Función auxiliar para generar nombres de archivo
            def generar_nombre_archivo(titulo, extension):
                """
                Genera un nombre de archivo válido basado en el título
                """
                nombre = titulo.lower()
                nombre = nombre.replace('á', 'a').replace('é', 'e').replace('í', 'i').replace('ó', 'o').replace('ú', 'u')
                nombre = ''.join(c for c in nombre if c.isalnum() or c.isspace())
                nombre = nombre.replace(' ', '_')[:50]
                return f"{nombre}.{extension}"

            # Obtener título del documento
            titulo_documento = lineas[0].strip()
            
            # 33. Procesamiento específico para Word
            if formato == 'docx':
                filename = generar_nombre_archivo(titulo_documento, 'docx')
                doc = Document()
                
                # Agregar título
                titulo_parrafo = doc.add_paragraph()
                run = titulo_parrafo.add_run(titulo_documento)
                run.bold = True
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0, 51, 102)
                titulo_parrafo.alignment = WD_ALIGN_PARAGRAPH.CENTER
                
                # Procesar contenido
                for linea in lineas[1:]:
                    if linea.strip():
                        parrafo = doc.add_paragraph()
                        if ':' in linea and len(linea) < 40:  # Subtítulos
                            run = parrafo.add_run(linea)
                            run.bold = True
                            run.font.size = Pt(14)
                            run.font.color.rgb = RGBColor(0, 51, 102)
                        else:  # Texto normal
                            run = parrafo.add_run(linea)
                            run.font.size = Pt(12)
                        
                        parrafo.spacing_after = Pt(12)
                
                doc.save(filename)
                
            # 34. Procesamiento específico para PDF
            elif formato == 'pdf':
                filename = generar_nombre_archivo(titulo_documento, 'pdf')
                c = canvas.Canvas(filename, pagesize=letter)
                width, height = letter
                
                # Configuración de márgenes
                margen_x = 72
                margen_y = 72
                y = height - margen_y
                ancho_util = width - (2 * margen_x)
                
                # 35. Función auxiliar para dibujar texto con ajuste de línea
                def draw_wrapped_line(texto, x, y, width, font_name, font_size):
                    """
                    Dibuja texto con ajuste de línea automático
                    """
                    c.setFont(font_name, font_size)
                    words = texto.split()
                    lines = []
                    current_line = []
                    current_width = 0

                    # Dividir texto en líneas
                    for word in words:
                        word_width = c.stringWidth(word, font_name, font_size)
                        space_width = c.stringWidth(' ', font_name, font_size)
                        
                        if current_width + word_width + space_width <= width:
                            current_line.append(word)
                            current_width += word_width + space_width
                        else:
                            lines.append(' '.join(current_line))
                            current_line = [word]
                            current_width = word_width + space_width
                    
                    if current_line:
                        lines.append(' '.join(current_line))

                    # Dibujar líneas
                    current_y = y
                    for line in lines:
                        if font_name == 'Helvetica-Bold' and font_size == 18:
                            text_width = c.stringWidth(line, font_name, font_size)
                            x_centered = (width - text_width) / 2
                            c.drawString(x_centered, current_y, line)
                        else:
                            c.drawString(x, current_y, line)
                        current_y -= font_size * 1.2

                    return current_y

                # 36. Función auxiliar para centrar título
                def center_title(texto, y, font_size=18):
                    """
                    Centra y formatea el título del documento
                    """
                    c.setFont('Helvetica-Bold', font_size)
                    c.setFillColorRGB(0, 0.2, 0.4)
                    
                    palabras = texto.split()
                    lineas = []
                    linea_actual = []
                    
                    # Dividir título en líneas
                    for palabra in palabras:
                        linea_actual.append(palabra)
                        if c.stringWidth(' '.join(linea_actual), 'Helvetica-Bold', font_size) > ancho_util * 0.8:
                            if len(linea_actual) > 1:
                                linea_actual.pop()
                                lineas.append(' '.join(linea_actual))
                                linea_actual = [palabra]
                            else:
                                lineas.append(' '.join(linea_actual))
                                linea_actual = []
                    
                    if linea_actual:
                        lineas.append(' '.join(linea_actual))
                    
                    # Dibujar líneas centradas
                    current_y = y
                    for linea in lineas:
                        text_width = c.stringWidth(linea, 'Helvetica-Bold', font_size)
                        x_centered = (width - text_width) / 2
                        c.drawString(x_centered, current_y, linea)
                        current_y -= font_size * 1.5
                    
                    return current_y - 20

                # 37. Función auxiliar para agregar números de página
                def add_page_number(canvas):
                    """
                    Añade número de página centrado al pie de página
                    """
                    canvas.saveState()
                    canvas.setFont('Helvetica', 10)
                    canvas.setFillColorRGB(0, 0, 0)
                    page_num = canvas.getPageNumber()
                    text = f"Página {page_num}"
                    text_width = canvas.stringWidth(text, 'Helvetica', 10)
                    x = (width - text_width) / 2
                    canvas.drawString(x, 30, text)
                    canvas.restoreState()

                # 38. Procesar contenido línea por línea para PDF
                for i, linea in enumerate(lineas):
                    if linea.strip():
                        if i == 0:  # Título
                            y = center_title(linea, y)
                        elif ':' in linea and len(linea) < 40:  # Subtítulos
                            c.setFillColorRGB(0, 0.2, 0.4)
                            y -= 10
                            y = draw_wrapped_line(linea, margen_x, y, ancho_util, 'Helvetica-Bold', 14)
                            y -= 15
                        else:  # Texto normal
                            c.setFillColorRGB(0, 0, 0)
                            y = draw_wrapped_line(linea, margen_x, y, ancho_util, 'Helvetica', 12)
                            y -= 12

                        # Nueva página si necesario
                        if y < margen_y + 50:
                            add_page_number(c)
                            c.showPage()
                            y = height - margen_y
                    else:
                        y -= 12

                # Agregar número a última página
                add_page_number(c)
                c.save()
            
            # 39. Procesamiento específico para TXT
            else:
                filename = generar_nombre_archivo(titulo_documento, 'txt')
                with open(filename, 'w', encoding='utf-8') as f:
                    f.write(contenido)
            
            return send_file(filename, as_attachment=True)
            
        except Exception as e:
            return f"Error al procesar la solicitud: {str(e)}"
    
    return render_template('index.html')

# 40. Ruta para apagar el servidor
@app.route('/shutdown', methods=['POST'])
def shutdown():
    """
    Ruta para apagar el servidor Flask de forma segura
    """
    func = request.environ.get('werkzeug.server.shutdown')
    if func is None:
        return 'El servidor no se puede apagar'
    func()
    return 'Servidor apagado'

# 41. Punto de entrada de la aplicación
if __name__ == '__main__':
    app.run(debug=True)